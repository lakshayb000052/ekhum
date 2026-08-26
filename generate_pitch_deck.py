import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation with 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Premium Light Mode Theme Color Palette
COLOR_BG_LIGHT = RGBColor(248, 250, 252)       # #F8FAFC (Soft Light Slate Background)
COLOR_WHITE = RGBColor(255, 255, 255)          # #FFFFFF (Pure Card White)
COLOR_BORDER_LIGHT = RGBColor(226, 232, 240)   # #E2E8F0 (Subtle Border)
COLOR_NAVY_TITLE = RGBColor(15, 23, 42)        # #0F172A (Deep Space Navy for Titles)
COLOR_TEXT_BODY = RGBColor(71, 85, 105)        # #475569 (Slate for readable body text)
COLOR_TEXT_MUTED = RGBColor(100, 116, 139)     # #64748B (Secondary Muted Text)
COLOR_EMERALD = RGBColor(5, 150, 105)          # #059669 (Primary Emerald Brand)
COLOR_EMERALD_LIGHT = RGBColor(16, 185, 129)   # #10B981 (Bright Emerald)
COLOR_TEAL = RGBColor(13, 148, 136)            # #0D9488 (Teal Accent)
COLOR_GOLD = RGBColor(217, 119, 6)             # #D97706 (Amber Accent)
COLOR_RED = RGBColor(220, 38, 38)              # #DC2626 (Alert Red)
COLOR_CARD_SUBTLE = RGBColor(241, 245, 249)    # #F1F5F9 (Subtle fill)

# Image paths (Light Mode Assets)
LOGO_LIGHT_PATH = r"C:\Users\laksh\.gemini\antigravity\brain\8b7bb95f-444d-4f03-89fc-79e092317602\ekhum_chaos_logo_light_1787079701569.jpg"
HERO_LIGHT_PATH = r"C:\Users\laksh\.gemini\antigravity\brain\8b7bb95f-444d-4f03-89fc-79e092317602\childfund_pitch_hero_light_1787079751707.jpg"

blank_slide_layout = prs.slide_layouts[6]

# Helper function to create background color
def set_slide_background(slide, color=COLOR_BG_LIGHT):
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    return bg_shape

# Helper function to create card shape
def add_card(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_BORDER_LIGHT, border_width=1):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    return card

# Helper for Header on Light slides
def add_slide_header(slide, category, title, subtitle=None):
    # Top Emerald Brand Strip
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.35), Inches(0.4), Inches(0.06))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_EMERALD
    strip.line.fill.background()

    # Category tag
    tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_EMERALD

    # Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_NAVY_TITLE

    if subtitle:
        p_sub = tf_title.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED

# ==============================================================================
# SLIDE 1: TITLE SLIDE (Light Mode Executive Presentation)
# ==============================================================================
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide1, COLOR_WHITE)

# Left emerald accent bar
accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), prs.slide_height)
accent.fill.solid()
accent.fill.fore_color.rgb = COLOR_EMERALD
accent.line.fill.background()

# Title text box
tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(7.5), Inches(4.8))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_badge = tf1.paragraphs[0]
p_badge.text = "CHAOS DESIGN  •  INDIVIDUAL GIVING FUNDRAISING SUITE"
p_badge.font.size = Pt(12)
p_badge.font.bold = True
p_badge.font.color.rgb = COLOR_EMERALD
p_badge.space_after = Pt(16)

p_main = tf1.add_paragraph()
p_main.text = "Scaling Child Sponsorship &\nIndividual Giving in India"
p_main.font.size = Pt(36)
p_main.font.bold = True
p_main.font.color.rgb = COLOR_NAVY_TITLE
p_main.space_after = Pt(16)

p_sub = tf1.add_paragraph()
p_sub.text = "An Enterprise Digital Giving Platform Purpose-Built Exclusively for ChildFund"
p_sub.font.size = Pt(16)
p_sub.font.color.rgb = COLOR_TEXT_BODY
p_sub.space_after = Pt(24)

p_footer = tf1.add_paragraph()
p_footer.text = "Presented by: EKhum Platform Team  |  A CHAOS Design Product\nMulti-Gateway  •  Automated Journeys  •  Section 80G & 10BD Statutory Engine"
p_footer.font.size = Pt(12)
p_footer.font.bold = True
p_footer.font.color.rgb = COLOR_EMERALD

# Add Light Mode Logo Image on Right
if os.path.exists(LOGO_LIGHT_PATH):
    slide1.shapes.add_picture(LOGO_LIGHT_PATH, Inches(8.8), Inches(1.5), width=Inches(3.8), height=Inches(3.8))

# ==============================================================================
# SLIDE 2: EXECUTIVE SUMMARY & VISION (Light Mode)
# ==============================================================================
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide2, COLOR_BG_LIGHT)
add_slide_header(slide2, "Executive Overview", "Transforming ChildFund's Digital Giving Ecosystem", "How EKhum accelerates child sponsorship acquisition, retention, and donor trust.")

cards_data2 = [
    ("⚡ Mission-Critical Reliability", "Multi-Gateway Redundancy", "Integrated with Razorpay, PayU, CCAvenue, and Worldline for zero transaction drop-offs and automated recurring retry mechanisms."),
    ("🎯 Automated Donor Engagement", "WhatsApp & Email Journeys", "Visual drag-and-drop lifecycle builder delivering child impact updates, birthday greetings, and automated mandate failure recovery."),
    ("📜 Frictionless Tax Compliance", "Instant 80G & Form 10BD", "Automated real-time 80G receipt PDF delivery via WhatsApp & Email + 1-click annual Form 10BD export for Income Tax compliance.")
]

for idx, (tag, title, desc) in enumerate(cards_data2):
    left = Inches(0.8 + idx * 4.0)
    top = Inches(1.8)
    width = Inches(3.7)
    height = Inches(4.8)
    
    add_card(slide2, left, top, width, height, COLOR_WHITE, COLOR_EMERALD, 1.5)
    
    tb = slide2.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), width - Inches(0.5), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = tag
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(12)
    
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 3: CHILDFUND KEY FUNDRAISING CHALLENGES (Light Mode)
# ==============================================================================
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide3, COLOR_BG_LIGHT)
add_slide_header(slide3, "The Opportunity", "Addressing Key Bottlenecks in Child Sponsorship Giving", "Analyzing friction points in traditional nonprofit donor management.")

challenges = [
    ("🛑 Mandate Attrition & Payment Drop-offs", "UPI Autopay and ENACH mandate failures often go unrecovered due to delayed manual outreach, resulting in a 25-35% drop in recurring child sponsorship retention."),
    ("💸 Single Gateway Vulnerability", "Dependency on a single payment aggregator causes critical revenue loss during gateway downtimes or bank maintenance windows during peak giving campaigns."),
    ("⏳ 80G Tax Support Overload", "Manual generation and dispatch of Section 80G certificates overwhelm support staff during March tax season, leading to delayed receipts and donor frustration."),
    ("📉 Disconnected Communication Channels", "Donors rarely receive continuous, transparent visual updates on the child they sponsor, leading to emotional detachment and lapsed pledges.")
]

for idx, (title, desc) in enumerate(challenges):
    row = idx // 2
    col = idx % 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(1.8 + row * 2.5)
    width = Inches(5.7)
    height = Inches(2.2)
    
    add_card(slide3, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide3.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(15)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_RED
    p0.space_after = Pt(8)
    
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(12)
    p1.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 4: THE SOLUTION (With Light Mode Hero Graphic)
# ==============================================================================
slide4 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide4, COLOR_BG_LIGHT)
add_slide_header(slide4, "The Solution", "EKhum by CHAOS: Purpose-Built for Indian Giving", "An integrated, cloud-native fundraising operating system with zero platform fees.")

features4 = [
    ("💳 4-Way Multi-Gateway Hub", "Razorpay, PayU, CCAvenue & Worldline routing with auto-fallback and unified reconciliation."),
    ("🤖 Visual Journey Builder", "Drag-and-drop workflow canvas for multi-step WhatsApp and Email donor journeys."),
    ("⚡ Real-Time Event Bus", "11 instant trigger events connecting donor actions to automated lifecycle responses."),
    ("📜 Automated 80G & 10BD Engine", "Instant statutory receipts with QR verification snapshot and Income Tax Form 10BD export.")
]

for idx, (title, desc) in enumerate(features4):
    row = idx // 2
    col = idx % 2
    left = Inches(0.8 + col * 3.4)
    top = Inches(1.8 + row * 2.5)
    width = Inches(3.2)
    height = Inches(2.3)
    
    add_card(slide4, left, top, width, height, COLOR_WHITE, COLOR_EMERALD, 1)
    
    tb = slide4.shapes.add_textbox(left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), height - Inches(0.36))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_NAVY_TITLE
    p0.space_after = Pt(6)
    
    p1 = tf.add_paragraph()
    p1.text = desc
    p1.font.size = Pt(11)
    p1.font.color.rgb = COLOR_TEXT_BODY

# Add Light Mode Hero Graphic on Right of Slide 4
if os.path.exists(HERO_LIGHT_PATH):
    slide4.shapes.add_picture(HERO_LIGHT_PATH, Inches(7.8), Inches(1.8), width=Inches(4.7), height=Inches(4.8))

# ==============================================================================
# SLIDE 5: MULTI-GATEWAY INFRASTRUCTURE (Light Mode)
# ==============================================================================
slide5 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide5, COLOR_BG_LIGHT)
add_slide_header(slide5, "Payment Infrastructure", "Multi-Gateway Redundancy & Smart Reconciliation", "Eliminating payment failures across Razorpay, PayU, CCAvenue & Worldline.")

gateways = [
    ("Razorpay", "UPI Autopay, Cards, Netbanking", "Primary gateway for online donation checkouts with instant webhook event ingestion."),
    ("PayU India", "ENACH Mandates, Subscriptions", "Secondary failover routing for recurring child sponsorship debits with high bank success rates."),
    ("CCAvenue", "Multi-Bank E-Mandates, 50+ Banks", "Deep tier-2/3 Indian bank coverage ensuring every regional donor can set up automatic pledges."),
    ("Worldline", "Direct Debit Cards, POS & Offline", "Enterprise payment gateway processing large corporate giving and event-based donations.")
]

for idx, (gw_name, gw_features, gw_desc) in enumerate(gateways):
    left = Inches(0.8 + idx * 3.0)
    top = Inches(1.8)
    width = Inches(2.75)
    height = Inches(4.8)
    
    add_card(slide5, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.25), width - Inches(0.4), height - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = f"GATEWAY #{idx+1}"
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(8)
    
    p1 = tf.add_paragraph()
    p1.text = gw_name
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = gw_features
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_GOLD
    p2.space_after = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = gw_desc
    p3.font.size = Pt(11)
    p3.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 6: VISUAL JOURNEY BUILDER (Light Mode)
# ==============================================================================
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide6, COLOR_BG_LIGHT)
add_slide_header(slide6, "Engagement Engine", "Automated Visual Journey Builder for ChildFund", "Designing personalized, high-touch sponsor communication flows visually.")

journey_steps = [
    ("1. Entry Trigger", "🎯 Event Triggered", "Sponsor signs up for 'Educate a Girl Child' campaign via website."),
    ("2. Immediate Action", "📧 Email Welcome Kit", "Dispatches personalized welcome kit with child bio, photo, and 80G receipt."),
    ("3. Time Delay", "⏳ Wait 48 Hours", "System pauses to allow donor to read impact materials before next touchpoint."),
    ("4. Mobile Nudge", "💬 WhatsApp Video", "Sends a warm 30-second introduction message from the ChildFund field team."),
    ("5. Smart Split", "🔀 Condition Branch", "Checks if donation > ₹5,000 to route to Major Donor relationship officer.")
]

for idx, (step_num, step_name, step_desc) in enumerate(journey_steps):
    left = Inches(0.8 + idx * 2.4)
    top = Inches(1.8)
    width = Inches(2.2)
    height = Inches(4.8)
    
    add_card(slide6, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide6.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = step_num
    p0.font.size = Pt(10)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(6)
    
    p1 = tf.add_paragraph()
    p1.text = step_name
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = step_desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 7: REAL-TIME EVENT TRIGGER ENGINE (Light Mode)
# ==============================================================================
slide7 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide7, COLOR_BG_LIGHT)
add_slide_header(slide7, "Automation Architecture", "Real-Time Event Trigger Engine", "11 standard system events that power autonomous giving operations.")

events_list = [
    ("donation.completed", "Fires on successful payment; auto-issues 80G receipt and triggers thank-you series."),
    ("donation.failed", "Fires when card or UPI payment drops; immediately dispatches 1-click retry WhatsApp link."),
    ("subscription.created", "Fires on new recurring child sponsorship pledge; enrols sponsor in onboarding journey."),
    ("subscription.cancelled", "Fires when recurring pledge is cancelled; initiates win-back survey sequence."),
    ("mandate.failed", "Fires when monthly bank auto-debit bounces; initiates automated 3-day recovery workflow."),
    ("receipt.generated", "Fires when 80G certificate PDF is rendered; delivers instant PDF attachment to donor.")
]

for idx, (ev_name, ev_desc) in enumerate(events_list):
    row = idx // 3
    col = idx % 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.8 + row * 2.5)
    width = Inches(3.7)
    height = Inches(2.2)
    
    add_card(slide7, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide7.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = f"⚡ {ev_name}"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(8)
    
    p1 = tf.add_paragraph()
    p1.text = ev_desc
    p1.font.size = Pt(11)
    p1.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 8: AUTOMATED 80G TAX COMPLIANCE & FORM 10BD (Light Mode)
# ==============================================================================
slide8 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide8, COLOR_BG_LIGHT)
add_slide_header(slide8, "Statutory Compliance", "Section 80G Tax Certificates & Form 10BD Ready", "Zero manual tax receipt generation and 100% compliant filings for ChildFund.")

tax_cards = [
    ("📜 Instant 80G Receipts", "Automated PDF Generation", "Every donation automatically generates a verifiable Section 80G PDF with child sponsorship details, donor PAN, organization 80G URN, and authorized signatory snapshot."),
    ("📤 1-Click Form 10BD Filing", "Income Tax Portal Export", "One-click export strictly formatted according to the Income Tax Department's Form 10BD schema, completely eliminating manual data formatting and reconciliation errors."),
    ("🛡️ Immutable Audit Snapshots", "Cryptographic Record Lock", "Preserves immutable snapshots of donor address and PAN at the exact moment of giving, guaranteeing flawless compliance during annual statutory audits.")
]

for idx, (tag, title, desc) in enumerate(tax_cards):
    left = Inches(0.8 + idx * 4.0)
    top = Inches(1.8)
    width = Inches(3.7)
    height = Inches(4.8)
    
    add_card(slide8, left, top, width, height, COLOR_WHITE, COLOR_EMERALD, 1.5)
    
    tb = slide8.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), width - Inches(0.5), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = tag
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(12)
    
    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 9: CHILDFUND SPONSORSHIP CRM & DATA MODEL (Light Mode)
# ==============================================================================
slide9 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide9, COLOR_BG_LIGHT)
add_slide_header(slide9, "Donor Data Architecture", "Schema.xlsx Extended Data Architecture", "Built specifically around the nuances of Indian individual giving and sponsorship tracking.")

schema_cards = [
    ("👤 Contact Object (26 Fields)", "Full Indian KYC & Identity", "Validates 10-digit PAN (ABCDE1234F), 6-digit PIN codes, family relationship links, total lifetime gifts, and acquisition campaign tagging."),
    ("🔄 Monthly Donation Object (28 Fields)", "Recurring Sponsorship Tracking", "Manages recurring child sponsorship pledges, mandate IDs, consecutive failed installments, preferred debit days, and pledge value upgrade histories."),
    ("💳 Payment Object (12 Fields)", "Transaction Ledger & Reconciliation", "Captures raw gateway transaction IDs, settlement dates, bank UTRs, net settled amounts, payment methods (UPI/Card/ENACH), and 80G receipt mappings.")
]

for idx, (title, subtitle, desc) in enumerate(schema_cards):
    left = Inches(0.8 + idx * 4.0)
    top = Inches(1.8)
    width = Inches(3.7)
    height = Inches(4.8)
    
    add_card(slide9, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide9.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), width - Inches(0.5), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_NAVY_TITLE
    p0.space_after = Pt(8)
    
    p1 = tf.add_paragraph()
    p1.text = subtitle
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_EMERALD
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 10: IMPACT & ROI PROJECTIONS (Light Mode)
# ==============================================================================
slide10 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide10, COLOR_BG_LIGHT)
add_slide_header(slide10, "Business Case & Impact", "Projected ROI & Key Metric Improvements", "Quantifiable efficiency gains within the first 6 months of deployment.")

roi_metrics = [
    ("+35%", "Sponsor Retention", "Automated WhatsApp lifecycle nurturing and payment failure recovery sequences significantly boost renewal rates."),
    ("-70%", "Payment Drop-Offs", "4-way gateway redundancy and smart failover routing preserve transactions that would otherwise fail on a single provider."),
    ("0.0%", "Platform Fee", "100% of donor funding goes directly to ChildFund's field operations without platform commissions or revenue cuts."),
    ("85%", "Faster 80G Delivery", "Real-time automated PDF delivery reduces donor inquiries during March tax season to nearly zero.")
]

for idx, (val, label, desc) in enumerate(roi_metrics):
    left = Inches(0.8 + idx * 3.0)
    top = Inches(1.8)
    width = Inches(2.75)
    height = Inches(4.8)
    
    add_card(slide10, left, top, width, height, COLOR_WHITE, COLOR_EMERALD, 1.5)
    
    tb = slide10.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = val
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(6)
    
    p1 = tf.add_paragraph()
    p1.text = label
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 11: SECURITY, DATA SOVEREIGNTY & MULTITENANCY (Light Mode)
# ==============================================================================
slide11 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide11, COLOR_BG_LIGHT)
add_slide_header(slide11, "Governance & Security", "Enterprise Security & DPDP Act Data Compliance", "Ensuring ChildFund's donor records and vulnerable child data remain strictly protected.")

sec_cards = [
    ("🔒 Strict Multitenancy", "Complete Data Isolation", "Every database object is strictly segregated by Organization ID with parameterized SQL query isolation preventing any cross-tenant data leakage."),
    ("🛡️ DPDP Act Consent Engine", "Explicit Consent Registry", "Full capture of donor communication consents across Email, WhatsApp, and SMS with immutable timestamps and IP address logging."),
    ("🔑 Granular RBAC Permissions", "Role-Based Access Control", "Distinct permission layers for Superadmins, NGO Leadership, Campaign Managers, and Finance Auditors.")
]

for idx, (title, subtitle, desc) in enumerate(sec_cards):
    left = Inches(0.8 + idx * 4.0)
    top = Inches(1.8)
    width = Inches(3.7)
    height = Inches(4.8)
    
    add_card(slide11, left, top, width, height, COLOR_WHITE, COLOR_BORDER_LIGHT, 1)
    
    tb = slide11.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), width - Inches(0.5), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_NAVY_TITLE
    p0.space_after = Pt(6)
    
    p1 = tf.add_paragraph()
    p1.text = subtitle
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_EMERALD
    p1.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_TEXT_BODY

# ==============================================================================
# SLIDE 12: DEPLOYMENT ROADMAP (Light Mode)
# ==============================================================================
slide12 = prs.slides.add_slide(blank_slide_layout)
set_slide_background(slide12, COLOR_BG_LIGHT)
add_slide_header(slide12, "Deployment Plan", "4-Week Rapid Onboarding & Pilot Launch", "A phased transition plan designed for zero disruption to ongoing child sponsorships.")

phases = [
    ("Week 1: Setup", "Gateway & WhatsApp Setup", "Configure ChildFund sub-merchant accounts across Razorpay, PayU & Meta WhatsApp Business API."),
    ("Week 2: Migration", "Legacy Data Migration", "Import existing donor database, historical child sponsorship records, and active ENACH mandates."),
    ("Week 3: Journeys", "Workflow Customization", "Build customized Welcome, Child Impact Update, and Payment Recovery journey flows."),
    ("Week 4: Go-Live", "Pilot Campaign Launch", "Launch live child sponsorship giving page, enable 80G automation, and train ChildFund teams.")
]

for idx, (p_num, p_title, p_desc) in enumerate(phases):
    left = Inches(0.8 + idx * 3.0)
    top = Inches(1.8)
    width = Inches(2.75)
    height = Inches(4.8)
    
    add_card(slide12, left, top, width, height, COLOR_WHITE, COLOR_EMERALD, 1.5)
    
    tb = slide12.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), height - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = p_num
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_EMERALD
    p0.space_after = Pt(8)
    
    p1 = tf.add_paragraph()
    p1.text = p_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_NAVY_TITLE
    p1.space_after = Pt(12)
    
    p2 = tf.add_paragraph()
    p2.text = p_desc
    p2.font.size = Pt(11.5)
    p2.font.color.rgb = COLOR_TEXT_BODY

# Save presentation to both exact requested names: child_fund_pitch.ppt and child_fund_pitch.pptx
output_ppt_workspace = r"e:\DanaPro\child_fund_pitch.ppt"
output_pptx_workspace = r"e:\DanaPro\child_fund_pitch.pptx"

output_ppt_brain = r"C:\Users\laksh\.gemini\antigravity\brain\8b7bb95f-444d-4f03-89fc-79e092317602\child_fund_pitch.ppt"
output_pptx_brain = r"C:\Users\laksh\.gemini\antigravity\brain\8b7bb95f-444d-4f03-89fc-79e092317602\child_fund_pitch.pptx"

prs.save(output_ppt_workspace)
prs.save(output_pptx_workspace)
prs.save(output_ppt_brain)
prs.save(output_pptx_brain)

print(f"SUCCESS: Generated child_fund_pitch.ppt and child_fund_pitch.pptx at {output_ppt_workspace}")
